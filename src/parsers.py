"""
文档解析器 - 支持 Word、PDF、PPT 格式
支持图片 OCR 文字识别（使用 PaddleOCR）
"""

import os
import io
import base64
from pathlib import Path
from typing import List, Dict, Optional, Tuple
import traceback


class DocumentParser:
    """统一的文档解析接口"""

    def __init__(self, enable_ocr: bool = True, ocr_lang: str = 'ch'):
        """
        初始化文档解析器

        Args:
            enable_ocr: 是否启用图片 OCR 识别（默认 True）
            ocr_lang: OCR 语言，'ch'中文，'en'英文（默认 'ch'）
        """
        self.supported_formats = {
            '.docx': self._parse_word,
            '.doc': self._parse_doc,
            '.pdf': self._parse_pdf,
            '.pptx': self._parse_ppt,
        }
        self.enable_ocr = enable_ocr
        self.ocr_lang = ocr_lang
        self._ocr_engine = None  # 延迟初始化

    def _get_ocr_engine(self):
        """获取 OCR 引擎（延迟初始化）"""
        if not self.enable_ocr:
            return None

        if self._ocr_engine is None:
            try:
                from paddleocr import PaddleOCR
                # 初始化 PaddleOCR - 强制使用CPU（避免与PyTorch GPU冲突）
                self._ocr_engine = PaddleOCR(
                    use_angle_cls=False,  # 禁用方向分类器以减少内存占用
                    lang=self.ocr_lang,
                    use_gpu=False,  # OCR强制使用CPU
                    show_log=False,  # 关闭日志输出
                    enable_mkldnn=False,  # 禁用MKL-DNN避免OneDnnContext兼容问题
                )
            except ImportError:
                print("警告: 未安装 PaddleOCR，图片文字识别功能不可用")
                print("      安装命令: pip install paddleocr paddlepaddle")
                self._ocr_engine = None
            except Exception as e:
                print(f"警告: PaddleOCR 初始化失败: {e}")
                self._ocr_engine = None

        return self._ocr_engine

    def _ocr_image(self, image_data: bytes, image_name: str = "image") -> str:
        """
        对图片进行 OCR 文字识别

        Args:
            image_data: 图片二进制数据
            image_name: 图片名称（用于日志）

        Returns:
            识别出的文本，失败返回空字符串
        """
        import gc

        ocr = self._get_ocr_engine()
        if ocr is None:
            return ""

        image = None
        img_array = None

        try:
            from PIL import Image
            import numpy as np

            # 将图片数据转换为 PIL Image
            image = Image.open(io.BytesIO(image_data))

            # 转换为 numpy 数组（RGB 格式）
            if image.mode != 'RGB':
                image = image.convert('RGB')
            img_array = np.array(image)

            # 使用 PaddleOCR 识别
            results = ocr.ocr(img_array, cls=False)  # 禁用分类器

            # 提取文字
            if results and results[0]:
                texts = []
                for line in results[0]:
                    if line and len(line) >= 2:
                        # line[0] 是坐标，line[1] 是 (文字, 置信度)
                        text = line[1][0] if line[1] else ""
                        if text and text.strip():
                            texts.append(text.strip())

                if texts:
                    ocr_text = "\n".join(texts)
                    return f"[图片文字 {image_name}]\n{ocr_text}\n"

        except Exception as e:
            print(f"图片 {image_name} OCR 识别失败: {e}")

        finally:
            # 显式清理资源
            if image is not None:
                image.close()
                del image
            if img_array is not None:
                del img_array
            gc.collect()

        return ""

    def parse(self, file_path: str, max_pages: int = None) -> Dict:
        """
        解析文档

        Args:
            file_path: 文档路径
            max_pages: 最大解析页数（None=全部，用于快速分类等场景）

        Returns:
            {
                'text': str,           # 提取的完整文本
                'pages': List[str],    # 按页/段落分割的文本
                'metadata': Dict       # 文档元数据
            }
        """
        file_path = Path(file_path)

        if not file_path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")

        ext = file_path.suffix.lower()

        if ext not in self.supported_formats:
            raise ValueError(f"不支持的文件格式: {ext}。支持的格式: {list(self.supported_formats.keys())}")

        parser_func = self.supported_formats[ext]

        try:
            # 传入 max_pages 参数
            result = parser_func(str(file_path), max_pages=max_pages)
            result['metadata'] = {
                'filename': file_path.name,
                'filepath': str(file_path),
                'format': ext,
                'size': file_path.stat().st_size
            }
            return result
        except Exception as e:
            raise Exception(f"解析文件 {file_path.name} 失败: {str(e)}\n{traceback.format_exc()}")

    def _parse_word(self, file_path: str, max_pages: int = None) -> Dict:
        """解析 Word (.docx) 文件 - 提取段落、表格、页眉、页脚、文本框、图片等"""
        try:
            from docx import Document
            from docx.oxml import parse_xml
            import zipfile
        except ImportError:
            raise ImportError("请安装 python-docx: pip install python-docx")

        # 验证文件是否是有效的 docx 格式（.docx 实际上是 ZIP 文件）
        try:
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                # 检查是否包含必要的 .docx 文件
                has_content = '[Content_Types].xml' in zip_ref.namelist()
                if not has_content:
                    raise Exception("文件不是有效的 .docx 格式（可能是 .doc 文件改了扩展名）")
        except zipfile.BadZipFile:
            raise Exception("文件已损坏或不是有效的 .docx 格式")
        except Exception as e:
            raise Exception(f"文件验证失败: {e}")

        try:
            doc = Document(file_path)
        except Exception as e:
            raise Exception(f"无法解析文件: {e}")

        all_text_parts = []

        # Word 文档没有明确的页数，使用段落数量估算
        # 假设每页约 25 个段落（可根据实际情况调整）
        estimated_pages = len(doc.paragraphs) // 25 + 1

        # 限制处理段落数量
        if max_pages and max_pages < estimated_pages:
            max_paragraphs = max_pages * 25
            print(f"  [DOC] Word 估算页数: {estimated_pages}，仅处理前 {max_pages} 页（约 {max_paragraphs} 段）")
        else:
            max_paragraphs = None
            print(f"  [DOC] Word 文档（共 {len(doc.paragraphs)} 段）")

        # 1. 提取页眉
        for section in doc.sections:
            if section.header:
                header = section.header
                for para in header.paragraphs:
                    text = para.text.strip()
                    if text:
                        all_text_parts.append(f"[页眉] {text}")

        # 2. 提取主段落文本（可能受 max_paragraphs 限制）
        paragraph_count = 0
        for para in doc.paragraphs:
            if max_paragraphs and paragraph_count >= max_paragraphs:
                break
            text = para.text.strip()
            if text:
                all_text_parts.append(text)
                paragraph_count += 1

        # 3. 提取表格文本（也受限制，避免处理过多表格）
        table_count = 0
        max_tables = max_paragraphs // 10 if max_paragraphs else None  # 假设每10段1个表格
        for table_idx, table in enumerate(doc.tables):
            if max_tables and table_count >= max_tables:
                break
            for row_idx, row in enumerate(table.rows):
                row_text = " | ".join([cell.text.strip() for cell in row.cells])
                if row_text.strip():
                    all_text_parts.append(row_text)
            table_count += 1

        # 4. 提取页脚
        for section in doc.sections:
            if section.footer:
                footer = section.footer
                for para in footer.paragraphs:
                    text = para.text.strip()
                    if text:
                        all_text_parts.append(f"[页脚] {text}")

        # 5. 提取文本框（从 XML 中）
        try:
            # 获取文档的 XML
            for para in doc.paragraphs:
                # 检查段落中是否包含文本框（通过检查 XML）
                if para._element.xml:
                    # 查找所有文本框（w:txbxContent 或 w:textbox）
                    xml = para._element.xml
                    if 'txbxContent' in xml or 'textbox' in xml:
                        # 尝试提取文本框内容
                        for child in para._element.iter():
                            if hasattr(child, 'text') and child.text and child.text.strip():
                                text = child.text.strip()
                                if text and text not in all_text_parts:
                                    all_text_parts.append(f"[文本框] {text}")
        except Exception as e:
            # 文本框提取失败不影响主流程
            pass

        # 6. 提取嵌入对象（如 Excel 表格等）
        try:
            from docx.oxml.ns import qn
            for para in doc.paragraphs:
                # 检查嵌入对象
                for obj in para._element.iter():
                    if obj.tag.endswith('}object'):  # 嵌入对象
                        # 尝试提取对象的数据
                        for child in obj.iter():
                            if hasattr(child, 'text') and child.text and child.text.strip():
                                text = child.text.strip()
                                if text and text not in all_text_parts:
                                    all_text_parts.append(f"[嵌入对象] {text}")
        except Exception as e:
            pass

        # ===== 7. 提取并识别图片中的文字 =====
        if self.enable_ocr:
            try:
                # 从文档中提取图片
                # Word 文档的图片存储在 document.xml.rels 中
                image_count = 0
                for rel in doc.part.rels.values():
                    if "image" in rel.target_ref:
                        try:
                            image_data = rel.target_part.blob
                            image_count += 1

                            # 使用 OCR 识别图片中的文字
                            ocr_text = self._ocr_image(
                                image_data,
                                image_name=f"word_img{image_count}"
                            )

                            if ocr_text:
                                all_text_parts.append(ocr_text)

                        except Exception as e:
                            # 单个图片识别失败不影响其他图片
                            continue
            except Exception as e:
                # 图片提取失败不影响文本解析
                pass

        # 合并所有文本
        full_text = "\n".join(all_text_parts)

        # 按段落分页（用于后续处理）
        paragraphs = [t for t in all_text_parts if not t.startswith('[页眉]') and not t.startswith('[页脚]')]
        pages = paragraphs if paragraphs else [full_text]

        return {
            'text': full_text,
            'pages': pages,
            'type': 'word'
        }

    def _parse_doc(self, file_path: str, max_pages: int = None) -> Dict:
        """解析旧版 Word (.doc) 文件"""
        # 方法1: 使用 pyth 库解析（推荐，不需要 Word）
        try:
            import pyth.plugins.rtf15.reader
            from pyth.rtf16.reader import Rtf16Reader
            from pyth.reader import Reader
            import io

            print(f"  [DOC] 使用 pyth 库解析...")

            # 读取文件
            with open(file_path, 'rb') as f:
                content = f.read()

            # 使用 pyth 解析
            reader = Reader(io.BytesIO(content))
            doc = reader.read()

            # 提取文本
            text = doc.plain_text()

            if not text or len(text.strip()) < 10:
                raise Exception("提取内容为空")

            # 如果指定了 max_pages，估算并截取文本
            if max_pages:
                estimated_chars = len(text)
                estimated_pages = estimated_chars // 2000 + 1
                max_chars = max_pages * 2000

                if estimated_chars > max_chars:
                    print(f"  [DOC] .doc 估算页数: {estimated_pages}，仅处理前 {max_pages} 页")
                    text = text[:max_chars] + "\n...[内容截止]"
                else:
                    print(f"  [DOC] .doc 文档（估算 {estimated_pages} 页，共 {len(text)} 字符）")
            else:
                print(f"  [DOC] .doc 文档（共 {len(text)} 字符）")

            # 按段落分割
            paragraphs = [p.strip() for p in text.split('\n') if p.strip() and len(p.strip()) > 5]
            pages = paragraphs if paragraphs else [text]

            return {
                'text': text,
                'pages': pages,
                'type': 'doc'
            }

        except ImportError:
            print("警告: 未安装 pyth")
            print("      安装命令: pip install pyth")
        except Exception as e:
            print(f"pyth 解析失败: {e}，尝试其他方法...")

        # 方法2: 使用 doc2docx + pywin32（需要 Word）
        try:
            import doc2docx
            import tempfile
            import os

            print(f"  [DOC] 使用 doc2docx 转换为 .docx...")

            # 创建临时文件
            with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
                tmp_docx_path = tmp.name

            try:
                # 转换文件 (doc2docx 需要安装 Word)
                doc2docx.convert(str(file_path), tmp_docx_path)

                # 使用 _parse_word 解析转换后的文件
                result = self._parse_word(tmp_docx_path, max_pages=max_pages)
                result['type'] = 'doc'

                print(f"  [DOC] 转换成功")

                return result
            finally:
                # 删除临时文件
                if os.path.exists(tmp_docx_path):
                    os.unlink(tmp_docx_path)

        except ImportError:
            print("警告: 未安装 doc2docx")
            print("      安装命令: pip install doc2docx")
        except Exception as e:
            print(f"doc2docx 转换失败: {e}，尝试其他方法...")

        # 方法3: 使用 win32com（需要 Word）
        try:
            import win32com
            from win32com.client import Dispatch
            import tempfile
            import os

            print(f"  [DOC] 使用 Word COM 接口转换...")

            # 创建临时文件
            with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
                tmp_docx_path = tmp.name

            try:
                # 使用 Word COM 接口打开并另存为 .docx
                word = Dispatch('Word.Application')
                word.Visible = False

                doc = word.Documents.Open(str(file_path))
                doc.SaveAs2(tmp_docx_path, FileFormat=16)  # 16 = wdFormatXMLDocument
                doc.Close()
                word.Quit()

                # 使用 _parse_word 解析转换后的文件
                result = self._parse_word(tmp_docx_path, max_pages=max_pages)
                result['type'] = 'doc'

                print(f"  [DOC] COM 转换成功")

                return result
            finally:
                # 确保退出 Word
                try:
                    word.Quit()
                except:
                    pass
                # 删除临时文件
                if os.path.exists(tmp_docx_path):
                    os.unlink(tmp_docx_path)

        except ImportError:
            print("警告: 未安装 pywin32")
            print("      安装命令: pip install pywin32")
        except Exception as e:
            print(f"win32com 转换失败: {e}")

        # 所有方法都失败，抛出异常
        raise ImportError(
            "无法解析 .doc 文件。请安装以下任一库：\n"
            "  1. pip install pyth (推荐，纯 Python，不需要 Word)\n"
            "  2. pip install doc2docx (需要安装 Microsoft Word)\n"
            "  3. pip install pywin32 (仅 Windows，需要安装 Microsoft Word)\n\n"
            "提示：如果没有 Word，建议手动将 .doc 转为 .docx 或使用 LibreOffice 批量转换"
        )

    def _parse_pdf(self, file_path: str, max_pages: int = None) -> Dict:
        """解析 PDF 文件 - 支持中文，多方法回退"""
        # 方法1: 尝试使用 PyMuPDF (fitz) - 对中文支持最好
        try:
            import fitz  # PyMuPDF
            return self._parse_pdf_with_pymupdf(file_path, max_pages=max_pages)
        except ImportError:
            print("警告: 未安装 PyMuPDF，尝试使用 pdfplumber (pip install pymupdf 以获得更好的中文支持)")
        except Exception as e:
            print(f"PyMuPDF 解析失败: {e}，尝试使用 pdfplumber...")

        # 方法2: 回退到 pdfplumber，使用 layout 参数
        try:
            import pdfplumber
            return self._parse_pdf_with_pdfplumber(file_path, max_pages=max_pages)
        except ImportError:
            raise ImportError("请安装 PDF 解析库: pip install pymupdf pdfplumber")
        except Exception as e:
            # 方法3: 最后尝试使用 PyPDF2
            try:
                import PyPDF2
                return self._parse_pdf_with_pypdf2(file_path, max_pages=max_pages)
            except ImportError:
                raise ImportError("请安装 PDF 解析库: pip install pymupdf 或 pip install pdfplumber")
            except Exception as e2:
                raise Exception(f"所有 PDF 解析方法都失败了: pdfplumber({e}), PyPDF2({e2})")

    def _parse_pdf_with_pymupdf(self, file_path: str, max_pages: int = None) -> Dict:
        """使用 PyMuPDF 解析 PDF - 对中文支持最好，支持图片 OCR"""
        import fitz  # PyMuPDF

        doc = fitz.open(file_path)
        total_pages = len(doc)

        # 限制解析页数
        if max_pages and max_pages < total_pages:
            pages_to_process = max_pages
            print(f"  [PDF] PDF 总页数: {total_pages}，仅解析前 {pages_to_process} 页")
        else:
            pages_to_process = total_pages
            print(f"  [PDF] PDF 总页数: {total_pages}")

        pages_text = []
        full_text_parts = []
        use_ocr_for_all = False  # 检测到乱码后，对后续所有页面使用OCR

        # 检测第一页是否有乱码，决定是否需要全程OCR
        if total_pages > 0:
            first_page = doc[0]
            first_text = first_page.get_text()
            if self._has_garbled_text(first_text):
                print(f"  [WARNING]  检测到PDF存在编码问题，将使用OCR识别所有页面...")
                use_ocr_for_all = True

        for page_num, page in enumerate(doc):
            # 如果设置了 max_pages，只处理前 N 页
            if max_pages and page_num >= max_pages:
                break
            print(f"  [READ] 处理第 {page_num + 1}/{total_pages} 页...", end="", flush=True)
            page_text = ""

            # 判断是否需要使用页面级OCR
            use_ocr_this_page = use_ocr_for_all
            text = page.get_text("text", flags=fitz.TEXT_PRESERVE_WHITESPACE | fitz.TEXT_PRESERVE_LIGATURES)
            if text and text.strip():
                text = text.strip()
                # 检测是否是乱码
                if self._is_valid_chinese_text(text):
                    page_text = text
                    print(f" [OK] 文本提取成功 ({len(text)} 字符)")
                else:
                    # 文本层乱码，标记需要整页OCR，跳过 blocks/dict/raw 回退
                    # （这些回退方式提取出来的也是同样的乱码文本）
                    use_ocr_this_page = True

            # 如果标记了需要整页OCR，直接跳到OCR分支，不做其他文本提取
            if not use_ocr_this_page and not page_text:
                # 方法4: 尝试原始模式（不带 flags）
                raw_text = page.get_text("text")
                if raw_text and raw_text.strip():
                    page_text = raw_text.strip()

            # ===== 提取并识别图片中的文字 =====
            ocr_used = False
            if self.enable_ocr and not use_ocr_this_page:
                try:
                    # 获取页面中的图片列表
                    image_list = page.get_images()

                    # 策略1: 如果有嵌入图片对象，识别图片中的文字
                    if image_list:
                        print(f" [SCAN] 发现 {len(image_list)} 张图片，使用 OCR 识别...", end="", flush=True)
                        ocr_results = []
                        for img_index, img in enumerate(image_list):
                            try:
                                # 获取图片的 xref
                                xref = img[0]

                                # 提取图片
                                base_image = doc.extract_image(xref)
                                if base_image and "image" in base_image:
                                    image_bytes = base_image["image"]
                                    image_ext = base_image.get("ext", "png")

                                    # 使用 OCR 识别图片中的文字
                                    ocr_text = self._ocr_image(
                                        image_bytes,
                                        image_name=f"page{page_num + 1}_img{img_index + 1}"
                                    )

                                    if ocr_text:
                                        ocr_results.append(ocr_text)

                            except Exception as e:
                                # 单个图片识别失败不影响其他图片
                                continue

                        if ocr_results:
                            page_text += "\n" + "\n".join(ocr_results)
                            print(f" [OK] OCR 识别成功")
                            ocr_used = True
                        else:
                            print(f" [WARNING]  OCR 未能识别文字")

                except Exception as e:
                    # 图片提取失败不影响文本解析
                    import traceback
                    traceback.print_exc()
                    pass

            # 策略2: 如果文本层为空/乱码，对整页进行OCR（扫描型PDF）
            if (use_ocr_this_page or not page_text or len(page_text.strip()) < 10) and self.enable_ocr:
                try:
                    print(f" [SCAN] 文本层为空，对整页进行 OCR...", end="", flush=True)
                    ocr_text = self._ocr_pdf_page(page, page_num + 1, total_pages)
                    if ocr_text:
                        page_text = ocr_text
                        ocr_used = True
                except Exception as e:
                    import traceback
                    traceback.print_exc()
                    pass

            if page_text:
                pages_text.append(page_text)
                full_text_parts.append(page_text)
                if not ocr_used and not use_ocr_this_page:
                    print(f" [OK] 完成")

            # 🔧 内存管理：每处理10页后清理一次内存（仅对大文件）
            if total_pages > 50 and (page_num + 1) % 10 == 0:
                import gc
                gc.collect()

        doc.close()
        print(f"  [STAT] 解析完成: 共 {total_pages} 页, {len(full_text_parts)} 页有效")

        full_text = "\n\n".join(full_text_parts)

        return {
            'text': full_text,
            'pages': pages_text,
            'type': 'pdf'
        }

    def _parse_pdf_with_pdfplumber(self, file_path: str, max_pages: int = None) -> Dict:
        """使用 pdfplumber 解析 PDF"""
        import pdfplumber

        pages_text = []
        full_text_parts = []

        with pdfplumber.open(file_path) as pdf:
            total_pages = len(pdf.pages)

            # 限制解析页数
            if max_pages and max_pages < total_pages:
                pages_to_process = max_pages
                print(f"  [PDF] PDF 总页数: {total_pages}，仅解析前 {pages_to_process} 页")
            else:
                pages_to_process = total_pages
                print(f"  [PDF] PDF 总页数: {total_pages}")

            for page_num, page in enumerate(pdf.pages):
                # 如果设置了 max_pages，只处理前 N 页
                if max_pages and page_num >= max_pages:
                    break

                # 使用 layout 模式提取，对表格和复杂布局更好
                text = page.extract_text(layout=True)
                if not text or not text.strip():
                    # 回退到普通模式
                    text = page.extract_text()

                if text:
                    text = text.strip()
                    if self._is_valid_chinese_text(text):
                        pages_text.append(text)
                        full_text_parts.append(text)

        full_text = "\n\n".join(full_text_parts)

        return {
            'text': full_text,
            'pages': pages_text,
            'type': 'pdf'
        }

    def _parse_pdf_with_pypdf2(self, file_path: str, max_pages: int = None) -> Dict:
        """使用 PyPDF2 解析 PDF - 最后的回退选项"""
        import PyPDF2

        pages_text = []
        full_text_parts = []

        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            total_pages = len(reader.pages)

            # 限制解析页数
            if max_pages and max_pages < total_pages:
                pages_to_process = max_pages
                print(f"  [PDF] PDF 总页数: {total_pages}，仅解析前 {pages_to_process} 页")
            else:
                pages_to_process = total_pages
                print(f"  [PDF] PDF 总页数: {total_pages}")

            for page_num, page in enumerate(reader.pages):
                # 如果设置了 max_pages，只处理前 N 页
                if max_pages and page_num >= max_pages:
                    break
                text = page.extract_text()
                if text:
                    text = text.strip()
                    pages_text.append(text)
                    full_text_parts.append(text)

        full_text = "\n\n".join(full_text_parts)

        return {
            'text': full_text,
            'pages': pages_text,
            'type': 'pdf'
        }

    def _has_garbled_text(self, text: str) -> bool:
        """
        检测文本是否包含PDF编码乱码

        通过以下方式检测：
        1. 特定的乱码特征字符（犐、犆等PDF编码错误映射字符）
        2. 高位非标准Unicode字符（U+0080-U+00FF等 Latin-1 Supplement 区域）
           当此类字符占比过高时，通常是PDF字体映射错误

        Args:
            text: 待检测的文本

        Returns:
            True 表示检测到乱码，False 表示文本正常
        """
        if not text or len(text) < 10:
            return False

        # 常见的PDF编码错误导致的乱码字符
        garbled_chars = ['犐', '犆', '犛', '狊', '犠', '狅', '狀', '犻', '犵', '犳', '狔',
                         '犾', '犮', '犱', '狋', '狀', '犪', '犫', '狉', '犽', '犿',
                         '狀', '狆', 'q', '狊', '狋', '狌', '狏', '狑', '狓', '狔']

        # 统计乱码字符出现次数
        garbled_count = sum(1 for char in text if char in garbled_chars)

        text_len = len(text)

        # 如果乱码字符超过一定阈值（文本长度的0.5%或至少出现5次），认为是乱码
        threshold_ratio = 0.005  # 0.5%
        threshold_min = 5  # 至少5个乱码字符

        if garbled_count >= threshold_min or (garbled_count / text_len) > threshold_ratio:
            return True

        # 检测高位非标准Unicode字符（PDF字体映射乱码的典型特征）
        # 这类字符位于 Latin-1 Supplement (U+0080-U+00FF) 和 Latin Extended 区域
        # 正常中文PDF不应该大量出现这些字符
        high_byte_count = 0
        for char in text:
            code = ord(char)
            # 排除换行、制表符等控制字符
            if code <= 0x20:
                continue
            # 检测 0x80-0xFF 范围的高位字符（典型字体映射乱码特征）
            # 但排除常见的中文标点（如 0xA1-0xFF 的 GBK 高位字）
            # 这类字符在正常中文文本中不应大量出现
            if 0x80 <= code <= 0xFF:
                high_byte_count += 1

        # 如果高位字符占比超过 30%，很可能是乱码
        if text_len > 0 and high_byte_count / text_len > 0.30:
            return True

        return False

    def _ocr_pdf_page(self, page, page_num: int, total_pages: int) -> str:
        """
        使用OCR识别整个PDF页面

        Args:
            page: PyMuPDF页面对象
            page_num: 当前页码
            total_pages: 总页数

        Returns:
            识别出的文本
        """
        import gc
        import fitz  # PyMuPDF

        ocr = self._get_ocr_engine()
        if ocr is None:
            return ""

        pix = None
        img_data = None
        image = None
        img_array = None

        try:
            import numpy as np
            from PIL import Image

            # 渲染页面为图片（降低倍率以减少内存占用）
            mat = fitz.Matrix(1.5, 1.5)  # 从2倍降到1.5倍
            pix = page.get_pixmap(matrix=mat)
            img_data = pix.tobytes("png")

            # 转换为 PIL Image
            image = Image.open(io.BytesIO(img_data))

            # 转换为 RGB
            if image.mode != 'RGB':
                image = image.convert('RGB')

            # 转换为 numpy 数组
            img_array = np.array(image)

            # OCR 识别
            results = ocr.ocr(img_array, cls=False)  # 禁用分类器

            # 提取文字
            if results and results[0]:
                texts = []
                for line in results[0]:
                    if line and len(line) >= 2:
                        text = line[1][0] if line[1] else ""
                        if text and text.strip():
                            texts.append(text.strip())

                if texts:
                    ocr_text = "\n".join(texts)
                    print(f"     [OK] OCR识别成功 ({len(ocr_text)} 字符)")
                    return ocr_text

        except Exception as e:
            print(f"     [WARNING]  OCR识别失败: {e}")

        finally:
            # 显式清理资源
            if pix is not None:
                del pix
            if img_data is not None:
                del img_data
            if image is not None:
                image.close()
                del image
            if img_array is not None:
                del img_array
            gc.collect()

        return ""

    def _is_valid_chinese_text(self, text: str) -> bool:
        """检测文本是否包含有效的中文内容（非乱码）"""
        if not text or len(text) < 10:
            return True  # 太短的文本跳过检测

        # 首先检查PDF编码乱码
        if self._has_garbled_text(text):
            return False

        # 检查是否包含中文字符
        chinese_chars = sum(1 for c in text if '\u4e00' <= c <= '\u9fff')
        total_chars = len(text)
        if total_chars == 0:
            return False

        chinese_ratio = chinese_chars / total_chars

        # 如果中文字符占比超过 10%，认为是有效中文文本
        if chinese_ratio > 0.10:
            return True

        # 如果中文字符很少（< 10%），但文本主要是 ASCII 符号（非字母数字），
        # 很可能是 PDF 字体映射乱码（如 !"#$%&'()*+,-./0123456789）
        # 这类乱码的典型特征：大量 ASCII 标点符号 + 少量数字 + 少量中文水印
        if chinese_ratio < 0.10:
            import re
            # 统计 ASCII 非字母字符（纯符号：!"#$%&'()*+,-./ 等）
            # 排除空格、换行和 URL 中常见的字母
            non_letter_ascii = sum(1 for c in text
                                   if 0x21 <= ord(c) <= 0x7F
                                   and not c.isalpha()
                                   and c not in ' \n\r\t')
            non_letter_ratio = non_letter_ascii / total_chars
            # 如果纯符号字符占比过高（> 40%），很可能是字体映射乱码
            if non_letter_ratio > 0.40:
                # 进一步确认：检查是否有真正的中文词汇（连续2个以上中文字符）
                chinese_sequences = re.findall(r'[\u4e00-\u9fff]{2,}', text)
                meaningful_chinese = len(''.join(chinese_sequences))
                # 如果有意义的中文词汇很少（< 5%），确认为乱码
                if meaningful_chinese / total_chars < 0.05:
                    return False

        # 检查是否是乱码（大量重复的特殊字符）
        # 例如: "ÿÿþýüûÿþþýüýý"
        special_chars = set('ÿþüûýþ')
        special_count = sum(1 for c in text if c in special_chars)
        if special_count > total_chars * 0.3:
            return False  # 可能是乱码

        # 检查是否有基本的有效字符（字母、数字、中文、标点）
        valid_chars = sum(1 for c in text
                          if c.isalnum() or
                          '\u4e00' <= c <= '\u9fff' or
                          c in '，。！？、：；""''（）【】《》')
        if valid_chars < total_chars * 0.3:
            return False

        return True

    def _parse_ppt(self, file_path: str, max_pages: int = None) -> Dict:
        """解析 PowerPoint (.pptx) 文件 - 提取文本框、表格、SmartArt、图表、备注、图片等"""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("请安装 python-pptx: pip install python-pptx")

        prs = Presentation(file_path)
        total_slides = len(prs.slides)

        # 限制解析页数（对于PPT是幻灯片数）
        if max_pages and max_pages < total_slides:
            slides_to_process = max_pages
            print(f"  [STAT] PPT 总页数: {total_slides}，仅解析前 {slides_to_process} 页")
        else:
            slides_to_process = total_slides
            print(f"  [STAT] PPT 总页数: {total_slides}")

        slides_text = []
        full_text_parts = []

        # 用于存储 PPT 中的图片（延迟处理）
        ppt_images = []

        for slide_idx, slide in enumerate(prs.slides):
            # 如果设置了 max_pages，只处理前 N 页
            if max_pages and slide_idx >= max_pages:
                break
            slide_content = []

            # 提取所有形状中的文本
            for shape in slide.shapes:
                # 1. 普通文本框和占位符
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())

                # 2. 表格
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join([cell.text.strip() for cell in row.cells])
                        if row_text.strip():
                            slide_content.append(row_text)

                # 3. 组（Group）中的形状
                if shape.shape_type == 6:  # Group shape
                    try:
                        for sub_shape in shape.shapes:
                            if hasattr(sub_shape, "text") and sub_shape.text.strip():
                                slide_content.append(sub_shape.text.strip())
                    except:
                        pass

                # 4. SmartArt（作为图表处理）
                if shape.shape_type == 18:  # SmartArt
                    try:
                        # 尝试提取 SmartArt 中的文本
                        for node in shape.shapes:
                            if hasattr(node, "text") and node.text.strip():
                                slide_content.append(f"[SmartArt] {node.text.strip()}")
                    except:
                        # 如果无法访问子形状，尝试 XML 提取
                        try:
                            xml = shape.element.xml
                            # 查找所有文本节点
                            import re
                            texts = re.findall(r'<a:t>([^<]+)</a:t>', xml)
                            for text in texts:
                                if text.strip():
                                    slide_content.append(f"[SmartArt] {text.strip()}")
                        except:
                            pass

                # 5. 图表（Chart）
                if shape.has_chart:
                    try:
                        chart = shape.chart
                        # 提取图表标题
                        if chart.chart_title.has_text_frame:
                            title = chart.chart_title.text_frame.text.strip()
                            if title:
                                slide_content.append(f"[图表标题] {title}")

                        # 提取系列数据
                        for series in chart.series:
                            series_name = series.name if series.name else "系列"
                            slide_content.append(f"[图表系列] {series_name}")

                        # 提取类别（X轴标签）
                        if chart.category_collection:
                            categories = [cat.value if cat.value else "" for cat in chart.category_collection]
                            if categories:
                                slide_content.append(f"[图表类别] {', '.join(categories)}")
                    except Exception as e:
                        pass

                # 6. OLE 对象（嵌入对象）
                if shape.shape_type == 7:  # OLE Object
                    try:
                        # 尝试从 XML 中提取文本
                        xml = shape.element.xml
                        import re
                        texts = re.findall(r'<a:t>([^<]+)</a:t>', xml)
                        for text in texts:
                            if text.strip():
                                slide_content.append(f"[嵌入对象] {text.strip()}")
                    except:
                        pass

                # ===== 7. 提取图片（延迟 OCR 处理） =====
                if self.enable_ocr and shape.shape_type == 13:  # Picture
                    try:
                        # 获取图片数据
                        if hasattr(shape, 'image'):
                            image = shape.image
                            image_data = image.blob
                            image_filename = image.filename

                            # 存储图片信息，稍后统一处理
                            ppt_images.append({
                                'data': image_data,
                                'name': f"slide{slide_idx + 1}_{image_filename}",
                                'slide_idx': slide_idx
                            })
                    except Exception as e:
                        # 单个图片提取失败不影响其他图片
                        continue

            # 8. 提取备注
            try:
                if hasattr(slide, "notes_slide") and slide.notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_content.append(f"[备注] {notes_text}")
            except:
                pass

            # 合并幻灯片内容
            slide_text = "\n".join(slide_content)
            if slide_text:
                slides_text.append(slide_text)
                full_text_parts.append(slide_text)

        # ===== 处理所有图片的 OCR =====
        if self.enable_ocr and ppt_images:
            # 按幻灯片索引分组图片
            images_by_slide = {}
            for img in ppt_images:
                slide_idx = img['slide_idx']
                if slide_idx not in images_by_slide:
                    images_by_slide[slide_idx] = []
                images_by_slide[slide_idx].append(img)

            # 为每个幻灯片添加 OCR 文字
            for slide_idx, images in images_by_slide.items():
                ocr_texts = []
                for img in images:
                    ocr_text = self._ocr_image(img['data'], img['name'])
                    if ocr_text:
                        ocr_texts.append(ocr_text)

                # 将 OCR 文字添加到对应幻灯片的文本中
                if ocr_texts and slide_idx < len(full_text_parts):
                    full_text_parts[slide_idx] += "\n" + "\n".join(ocr_texts)

        full_text = "\n\n---\n\n".join(full_text_parts)

        return {
            'text': full_text,
            'pages': slides_text,
            'type': 'ppt'
        }

    def batch_parse(self, file_paths: List[str]) -> List[Dict]:
        """批量解析文档"""
        results = []
        for file_path in file_paths:
            try:
                result = self.parse(file_path)
                results.append(result)
            except Exception as e:
                print(f"警告: 解析文件 {file_path} 失败: {e}")
                continue
        return results


def scan_documents(directory: str, extensions: Optional[List[str]] = None) -> List[str]:
    """
    扫描目录中的所有支持格式的文档

    Args:
        directory: 目录路径
        extensions: 文件扩展名列表，如 ['.pdf', '.docx']

    Returns:
        文件路径列表
    """
    directory = Path(directory)

    if not directory.exists():
        raise FileNotFoundError(f"目录不存在: {directory}")

    if extensions is None:
        extensions = ['.pdf', '.docx', '.pptx', '.doc']

    file_paths = []

    for ext in extensions:
        file_paths.extend(directory.rglob(f"*{ext}"))

    return [str(p) for p in file_paths]


if __name__ == "__main__":
    # 测试代码
    parser = DocumentParser()

    # 测试单个文件解析
    # result = parser.parse("test.docx")
    # print(f"提取文本长度: {len(result['text'])}")
    # print(f"分段数量: {len(result['pages'])}")

    # 测试目录扫描
    # files = scan_documents("./documents")
    # print(f"找到 {len(files)} 个文档")
    # for f in files:
    #     print(f"  - {f}")
    pass
